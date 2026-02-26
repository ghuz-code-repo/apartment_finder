# app/services/presentation_service.py
import io
from datetime import date

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
except ImportError:
    plt = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    Presentation = None

GH_GOLD = RGBColor(0xC5, 0x95, 0x00)
GH_DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
GH_BLACK = RGBColor(0x10, 0x10, 0x10)
GH_GRAY = RGBColor(0x80, 0x80, 0x80)
LOGO_PATH = 'app/static/img/logo_gh.png'

def _to_hex(color_rgb):
    """Конвертация кортежа RGBColor в HEX."""
    return f"#{color_rgb[0]:02x}{color_rgb[1]:02x}{color_rgb[2]:02x}"

def _format_cell(cell, text, size=12, bold=False, align=PP_ALIGN.LEFT, color_rgb=None):
    if color_rgb is None: color_rgb = GH_BLACK
    tf = cell.text_frame; tf.clear(); p = tf.paragraphs[0]
    p.text = str(text); p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color_rgb
    p.alignment = align; tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

def _add_slide_footer(prs, slide):
    try: slide.shapes.add_picture(LOGO_PATH, Inches(15.0), Inches(8.3), height=Inches(0.6))
    except: pass
    box = slide.shapes.add_textbox(Inches(0.5), Inches(8.5), Inches(1.0), Inches(0.3))
    box.text_frame.text = f"Слайд {len(prs.slides)}"
    box.text_frame.paragraphs[0].font.size = Pt(10)

def _create_bar_chart_image(labels, data, title):
    if plt is None: raise ImportError("matplotlib не установлена.")
    fig, ax = plt.subplots(figsize=(14, 6))
    bars = ax.barh(labels, data, color=_to_hex(GH_GOLD))
    ax.set_title(title, fontsize=16, fontweight='bold', color=_to_hex(GH_DARK_BLUE))
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    m = max(data) if data else 1
    for b in bars:
        ax.text(b.get_width() + (m * 0.01), b.get_y() + b.get_height()/2, f'{b.get_width():,.0f}', va='center', ha='left', fontsize=12, fontweight='bold', color=_to_hex(GH_DARK_BLUE))
    fig.tight_layout(); buf = io.BytesIO(); fig.savefig(buf, format='png', dpi=100, transparent=True); plt.close(fig); buf.seek(0); return buf

def generate_pricelist_pptx(complex_name, prop_type, percent, stats):
    if Presentation is None: raise ImportError("python-pptx не установлена.")
    prs = Presentation(); prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    # Титульный
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = f"Анализ изменения цен: {complex_name}"
    s1.placeholders[1].text = f"Объекты: {prop_type}\nПовышение цены дна: {percent*100:+.2f}%\nДата: {date.today().strftime('%d.%m.%Y')}"
    _add_slide_footer(prs, s1)
    # Инфографика
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    img = _create_bar_chart_image(['Текущие остатки, $', 'Новые остатки, $'], [stats['final_totals_before'], stats['final_totals_after']], "Изменение общей стоимости свободного склада")
    s2.shapes.add_picture(img, Inches(1), Inches(1.5), width=Inches(14))
    _add_slide_footer(prs, s2)
    # Цены дна
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Сравнение ключевых цен дна ($/м²)"
    t = s3.shapes.add_table(4, 3, Inches(1), Inches(2), Inches(14), Inches(4)).table
    for i, h in enumerate(["Показатель", "Было", "Стало"]): _format_cell(t.cell(0, i), h, bold=True, color_rgb=GH_DARK_BLUE)
    f_b, f_a = stats['floor_prices_before'], stats['floor_prices_after']
    m_data = [("Мин. цена дна", min(f_b), min(f_a)), ("Сред. цена дна", sum(f_b)/len(f_b), sum(f_a)/len(f_a)), ("Макс. цена дна", max(f_b), max(f_a))]
    for i, (l, b, a) in enumerate(m_data, 1):
        _format_cell(t.cell(i, 0), l); _format_cell(t.cell(i, 1), f"{b:,.0f}"); _format_cell(t.cell(i, 2), f"{a:,.0f}", bold=True, color_rgb=GH_GOLD)
    _add_slide_footer(prs, s3); buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf

def generate_passport_pptx(data, usd_rate):
    # Код генерации паспорта проекта из предыдущих этапов (title, kpi, team, plan-fact, remainders, competitors)
    # ...
    pass